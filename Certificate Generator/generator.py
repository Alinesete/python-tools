import pandas as pd
import os
import shutil
from pptx import Presentation
import win32com.client

def generate_certificates(xlsx_file, template_path, user_input, output_dir):
    # Open the PowerPoint file and get the first slide
    df = pd.read_excel(xlsx_file)

    # Extract the values in the "Nomes" column and save them in a list
    nomes_blob = df['Names'].dropna().tolist()
    nomes = [nome.upper() for nome in nomes_blob]

    # Make the output directory if it doesn't exist
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)

    # Loop through each name and generate a customized PowerPoint file
    for name in nomes:
        # Open the template file and create a new presentation object based on it
        output_path = os.path.join(output_dir, f'{user_input} - {name}.pdf')
        output_path_pdf = os.path.abspath(output_path)
        output_path_pp = os.path.abspath(output_path.replace('.pdf','.pptx'))
        shutil.copy2(template_path, output_path_pp)
        pr = Presentation(output_path_pp)

        # Replace the "NOME" tag with the current name
        for slide in pr.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text') and 'NAME' in shape.text:
                    # Store the original formatting
                    original_text = shape.text
                    original_font_name = shape.text_frame.paragraphs[0].runs[0].font.name
                    original_font_size = shape.text_frame.paragraphs[0].runs[0].font.size

                    # Replace the text
                    shape.text = original_text.replace('NAME', name)

                    # Apply the original formatting
                    shape.text_frame.paragraphs[0].runs[0].font.name = original_font_name
                    shape.text_frame.paragraphs[0].runs[0].font.size = original_font_size

        # Save the updated file as a PDF
        pr.save(output_path_pp)

        # Save the PowerPoint file as a PDF using the COM interface and delete the pptx file at the end
        try:
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            deck = powerpoint.Presentations.Open(output_path_pp)
            deck.SaveAs(output_path_pdf, 32) # 32 is the value for saving as PDF
            deck.Close()
            powerpoint.Quit()
            os.remove(output_path_pp)
        except Exception as e:
            print(e)
        finally:
            # Clean up the PowerPoint object
            del deck
            del powerpoint